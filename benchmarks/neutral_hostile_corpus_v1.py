from __future__ import annotations

"""
CMPCT Neutral/Hostile Benchmark Corpus Generator v1

This generator creates deterministic, byte-stable, *valid* real-world-like files across
workloads that favor different compression strategies.  It intentionally includes both
friendly and hostile material, so a format cannot win merely by recognizing one special
container pattern.

Footnote: the corpus is synthetic for redistribution and repeatability, but the formats,
file relationships, entropy profiles, sizes, and directory shapes are chosen to resemble
ordinary developer, office, media, analytics, backup, ML, and workstation data.
"""

from pathlib import Path
import argparse, bz2, csv, gzip, hashlib, io, json, lzma, math, os, shutil, sqlite3, struct, subprocess, tarfile, textwrap, wave, zipfile
from datetime import datetime, timedelta, timezone

import numpy as np
from PIL import Image, ImageDraw, ImageFilter
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from docx import Document
from docx.shared import Inches
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches as PptInches

SEED = 0xC0DEC0DE
RNG = np.random.default_rng(SEED)

def reseed(tag: str) -> None:
    global RNG
    # Footnote: every workload owns an independent deterministic PRNG stream, so it can be
    # regenerated, removed, or reordered without silently changing the bytes of other corpora.
    salt = int.from_bytes(hashlib.sha256(tag.encode('utf-8')).digest()[:8], 'little')
    RNG = np.random.default_rng(SEED ^ salt)

ROOT = Path('CMPCT_Neutral_Hostile_Corpus_v1')

WORDS = (
    "account adaptive agent allocation analysis api archive audit batch benchmark build cache "
    "campaign client cluster code commit compression configuration container corpus customer data "
    "deployment design document engine event experiment feature finance gateway index integrity invoice "
    "latency log manifest market media memory metadata metric model operator package pipeline policy "
    "project recovery release report repository request response revision service session snapshot source "
    "storage stream system task telemetry test transaction update validation vector version workflow"
).split()


def reset_dir(p: Path) -> None:
    shutil.rmtree(p, ignore_errors=True)
    p.mkdir(parents=True, exist_ok=True)


def randbytes(n: int) -> bytes:
    # Footnote: deterministic pseudo-random bytes are deliberately used instead of os.urandom.
    # They are still incompressible for practical codecs, while allowing exact corpus regeneration.
    return RNG.integers(0, 256, size=n, dtype=np.uint8).tobytes()


def sentence(min_words=8, max_words=22) -> str:
    n = int(RNG.integers(min_words, max_words + 1))
    w = RNG.choice(WORDS, size=n, replace=True).tolist()
    w[0] = w[0].capitalize()
    return " ".join(w) + "."


def make_photo(path: Path, size=(1920, 1080), quality=88, variant=0) -> None:
    """Create a photography-like JPEG with gradients, texture, blur, and hard edges."""
    w, h = size
    y, x = np.mgrid[0:h, 0:w]
    base = np.empty((h, w, 3), dtype=np.float32)
    base[..., 0] = 45 + 115 * (x / max(1, w-1)) + 28 * np.sin((y + variant*41) / 91)
    base[..., 1] = 50 + 105 * (y / max(1, h-1)) + 25 * np.cos((x + variant*67) / 113)
    base[..., 2] = 70 + 80 * ((x+y) / max(1, w+h-2))
    noise = RNG.normal(0, 12 + (variant % 4) * 2, size=(h, w, 1))
    arr = np.clip(base + noise, 0, 255).astype(np.uint8)
    im = Image.fromarray(arr, 'RGB').filter(ImageFilter.GaussianBlur(radius=0.35))
    d = ImageDraw.Draw(im)
    for i in range(18):
        x0 = int(RNG.integers(0, w)); y0 = int(RNG.integers(0, h))
        x1 = min(w, x0 + int(RNG.integers(30, 420))); y1 = min(h, y0 + int(RNG.integers(20, 260)))
        col = tuple(int(v) for v in RNG.integers(20, 235, size=3))
        d.rounded_rectangle((x0, y0, x1, y1), radius=12, outline=col, width=3)
    path.parent.mkdir(parents=True, exist_ok=True)
    im.save(path, format='JPEG', quality=quality, subsampling='4:2:0', optimize=True)


def make_png_ui(path: Path, size=(1600, 1000), variant=0) -> None:
    w, h = size
    im = Image.new('RGB', size, (246, 247, 249))
    d = ImageDraw.Draw(im)
    d.rectangle((0, 0, w, 70), fill=(31, 36, 48))
    for i in range(10):
        y = 105 + i * 78
        d.rounded_rectangle((55, y, w-55, y+55), radius=10, fill=(255,255,255), outline=(190,195,204))
        d.rectangle((80, y+17, 250 + (i*37+variant*17)%260, y+35), fill=(60,90,155))
        for j in range(4):
            x0 = 620 + j*180
            d.rectangle((x0, y+18, x0+95+(i*j)%45, y+34), fill=(135,142,154))
    # Footnote: PNG screenshots are structurally compressible but already DEFLATE-compressed,
    # which punishes archives that blindly re-compress encoded bytes.
    path.parent.mkdir(parents=True, exist_ok=True)
    im.save(path, format='PNG', optimize=True, compress_level=6)


def make_wav(path: Path, seconds=18, sr=44100, channels=2, variant=0) -> None:
    t = np.arange(sr * seconds, dtype=np.float64) / sr
    signal = 0.48*np.sin(2*np.pi*(220+variant*11)*t) + 0.22*np.sin(2*np.pi*(440+variant*7)*t)
    signal += 0.06*np.sin(2*np.pi*2.5*t)*np.sin(2*np.pi*880*t)
    signal += RNG.normal(0, 0.015, size=t.shape)
    pcm = np.clip(signal*32767, -32768, 32767).astype('<i2')
    stereo = np.column_stack([pcm, np.roll(pcm, 13 + variant)]).ravel().tobytes()
    path.parent.mkdir(parents=True, exist_ok=True)
    with wave.open(str(path), 'wb') as wf:
        wf.setnchannels(channels); wf.setsampwidth(2); wf.setframerate(sr); wf.writeframes(stereo)


def run_ffmpeg(args: list[str]) -> None:
    subprocess.run(['ffmpeg', '-hide_banner', '-loglevel', 'error', '-y', *args], check=True)


def tree_hash(root: Path) -> str:
    h = hashlib.sha256()
    for p in sorted(q for q in root.rglob('*') if q.is_file()):
        rel = p.relative_to(root).as_posix().encode()
        b = p.read_bytes()
        h.update(len(rel).to_bytes(4,'little')); h.update(rel)
        h.update(len(b).to_bytes(8,'little')); h.update(b)
    return h.hexdigest()


def corpus_source_repo(root: Path) -> None:
    reseed('corpus_source_repo')
    p = root/'01_developer_repository'; reset_dir(p)
    (p/'src').mkdir(); (p/'tests').mkdir(); (p/'docs').mkdir(); (p/'build').mkdir(); (p/'.git_sim'/'objects').mkdir(parents=True)
    templates = [
        "def {name}(items, config):\n    total = 0\n    for item in items:\n        if item.get('enabled', True):\n            total += item.get('value', 0)\n    return {{'name': '{name}', 'total': total, 'config': config}}\n",
        "export function {name}(rows, options) {{\n  const active = rows.filter(r => r.enabled !== false);\n  return active.map((r, i) => ({{...r, rank: i + 1, source: '{name}'}}));\n}}\n",
    ]
    for i in range(420):
        ext = '.py' if i % 3 else '.ts'
        name = f'module_{i:04d}'
        body = templates[0 if ext=='.py' else 1].format(name=name)
        comments = '\n'.join(('# '+sentence()) if ext=='.py' else ('// '+sentence()) for _ in range(8 + i%9))
        (p/'src'/f'{name}{ext}').write_text(comments+'\n'+body, encoding='utf-8')
    for i in range(140):
        (p/'tests'/f'test_{i:04d}.txt').write_text('\n'.join(sentence() for _ in range(18)), encoding='utf-8')
    lock = {'lockfileVersion':3,'packages':{}}
    for i in range(4200):
        lock['packages'][f'node_modules/pkg-{i}']={'version':f'{1+i%8}.{i%25}.{i%113}','resolved':f'https://registry.example/pkg-{i}.tgz','integrity':hashlib.sha256(f'pkg{i}'.encode()).hexdigest()}
    (p/'package-lock.json').write_text(json.dumps(lock,separators=(',',':')), encoding='utf-8')
    (p/'docs'/'architecture.md').write_text('\n\n'.join(f'## {i}. {sentence(3,7)}\n'+ '\n'.join(sentence() for _ in range(12)) for i in range(80)), encoding='utf-8')
    # Build two real ELF executables with similar but not identical source.
    for i in range(2):
        csrc=p/'build'/f'app{i}.c'
        values=','.join(str((j*17+i*13)%997) for j in range(25000))
        csrc.write_text(f'#include <stdio.h>\nstatic int t[]={{ {values} }};\nint main(){{long s=0;for(int i=0;i<25000;i++)s+=t[i];printf("%ld\\n",s);return 0;}}\n')
        subprocess.run(['gcc','-O2','-s',str(csrc),'-o',str(p/'build'/f'app{i}')],check=True)
    # Git object-like compressed blobs: intentionally already-compressed small objects.
    import zlib
    for i in range(700):
        raw=(f'blob {i}\0'+('\n'.join(sentence() for _ in range(6+i%8)))).encode()
        h=hashlib.sha1(raw).hexdigest(); d=p/'.git_sim'/'objects'/h[:2]; d.mkdir(exist_ok=True)
        (d/h[2:]).write_bytes(zlib.compress(raw, 6))


def corpus_office(root: Path) -> None:
    reseed('corpus_office')
    p=root/'02_office_workspace'; reset_dir(p); assets=p/'assets'; assets.mkdir()
    for i in range(8): make_photo(assets/f'photo_{i}.jpg', (1600,900), 84+i%4, i)
    for i in range(5): make_png_ui(assets/f'dashboard_{i}.png', (1400,900), i)
    # DOCX family with shared embedded assets but independently zipped packages.
    for doci in range(4):
        doc=Document(); doc.add_heading(f'Quarterly Operating Review {doci+1}',0)
        for sec in range(18):
            doc.add_heading(f'{sec+1}. {sentence(3,7)}',level=1)
            for _ in range(4): doc.add_paragraph(sentence(18,34))
            if sec in {2,7,12}: doc.add_picture(str(assets/f'photo_{(sec+doci)%8}.jpg'), width=Inches(5.7))
        doc.save(p/f'operating_review_v{doci+1}.docx')
    wb=Workbook(); ws=wb.active; ws.title='Transactions'; ws.append(['date','customer','region','sku','units','price','discount','status'])
    base=datetime(2025,1,1)
    for i in range(45000):
        ws.append([(base+timedelta(days=i%580)).date().isoformat(),f'CUST-{i%3100:05d}', ['N','S','E','W'][i%4],f'SKU-{i%780:04d}',int(1+i%17),round(4.5+(i%97)*1.17,2),[0,.05,.1,.15][i%4],['open','paid','paid','refunded'][i%4]])
    s=wb.create_sheet('Summary'); s.append(['metric','value']);
    for i in range(500): s.append([f'metric_{i}', f'=SUM(Transactions!E{2+i}:E{100+i})'])
    wb.save(p/'finance_model.xlsx')
    prs=Presentation()
    for i in range(22):
        sl=prs.slides.add_slide(prs.slide_layouts[5]); sl.shapes.title.text=f'{i+1}. {sentence(3,8)}'
        sl.shapes.add_picture(str(assets/f'photo_{i%8}.jpg'), PptInches(0.8), PptInches(1.4), width=PptInches(5.6))
    prs.save(p/'board_deck.pptx')
    # PDF with text and JPEGs; valid cross-format repeated subobjects emerge naturally.
    pdf=p/'client_report.pdf'; c=canvas.Canvas(str(pdf),pagesize=letter)
    for page in range(28):
        c.setFont('Helvetica-Bold',15); c.drawString(54,742,f'Client Performance Report — Page {page+1}')
        c.setFont('Helvetica',9); y=714
        for _ in range(12): c.drawString(54,y,sentence(10,17)[:105]); y-=15
        if page%4==0: c.drawImage(str(assets/f'photo_{page%8}.jpg'),54,285,width=500,height=281,preserveAspectRatio=True,mask='auto')
        c.showPage()
    c.save()


def corpus_media(root: Path) -> None:
    reseed('corpus_media')
    p=root/'03_media_library'; reset_dir(p)
    for i in range(14): make_photo(p/f'IMG_{i:04d}.jpg',(2200,1467),quality=82+(i%5)*3,variant=20+i)
    for i in range(10): make_png_ui(p/f'screen_{i:03d}.png',(1920,1080),variant=40+i)
    make_wav(p/'field_recording.wav', seconds=26, variant=3)
    run_ffmpeg(['-f','lavfi','-i','testsrc2=size=1280x720:rate=30','-f','lavfi','-i','sine=frequency=523:sample_rate=48000','-t','10','-c:v','libx264','-preset','medium','-crf','22','-pix_fmt','yuv420p','-c:a','aac','-b:a','128k',str(p/'clip_h264.mp4')])
    run_ffmpeg(['-i',str(p/'field_recording.wav'),'-c:a','flac',str(p/'field_recording.flac')])
    run_ffmpeg(['-i',str(p/'field_recording.wav'),'-c:a','libmp3lame','-q:a','3',str(p/'field_recording.mp3')])


def corpus_analytics(root: Path) -> None:
    reseed('corpus_analytics')
    p=root/'04_analytics_and_database'; reset_dir(p)
    # CSV + JSONL share semantic values but not byte layout.
    regions=['north','south','east','west']; products=[f'P{i:04d}' for i in range(1200)]
    csvp=p/'events.csv'; jsonp=p/'events.jsonl'
    with csvp.open('w',newline='',encoding='utf-8') as fc, jsonp.open('w',encoding='utf-8') as fj:
        w=csv.writer(fc); w.writerow(['ts','account','region','product','qty','revenue','ok','session'])
        start=datetime(2026,1,1,tzinfo=timezone.utc)
        for i in range(90000):
            rec={'ts':(start+timedelta(seconds=i*17)).isoformat(),'account':f'A{i%19000:05d}','region':regions[i%4],'product':products[(i*7)%len(products)],'qty':1+i%13,'revenue':round(2.75+(i%191)*0.91,2),'ok':i%23!=0,'session':hashlib.md5(f'{i//9}'.encode()).hexdigest()[:16]}
            w.writerow(rec.values()); fj.write(json.dumps(rec,separators=(',',':'))+'\n')
    db=p/'warehouse.sqlite'; con=sqlite3.connect(db); cur=con.cursor()
    cur.executescript('PRAGMA page_size=4096; CREATE TABLE events(id INTEGER PRIMARY KEY, account TEXT, region TEXT, product TEXT, qty INT, revenue REAL, payload BLOB); CREATE INDEX idx_acc ON events(account);')
    rows=[]
    for i in range(45000): rows.append((f'A{i%19000:05d}',regions[i%4],products[(i*7)%len(products)],1+i%13,2.75+(i%191)*0.91,randbytes(24) if i%11==0 else b'status=ok;source=api'))
    cur.executemany('INSERT INTO events(account,region,product,qty,revenue,payload) VALUES(?,?,?,?,?,?)',rows); con.commit(); con.execute('VACUUM'); con.close()
    # Numeric arrays: float data is only moderately compressible; npz is already compressed.
    a=(RNG.normal(0,1,size=(1200,800)).astype('<f4'))
    np.save(p/'features.npy',a)
    np.savez_compressed(p/'features_compressed.npz',features=a,labels=np.arange(1200,dtype=np.int32)%17)


def corpus_logs(root: Path) -> None:
    reseed('corpus_logs')
    p=root/'05_logs_and_telemetry'; reset_dir(p)
    base=datetime(2026,7,1,tzinfo=timezone.utc)
    rawfiles=[]
    for day in range(6):
        q=p/f'app-{day:02d}.log'; rawfiles.append(q)
        with q.open('w',encoding='utf-8') as f:
            for i in range(14000):
                ts=(base+timedelta(days=day,seconds=i*2)).isoformat()
                level=['INFO','INFO','INFO','WARN','DEBUG'][i%5]
                route=['/api/jobs','/api/files','/api/search','/health'][i%4]
                f.write(f'{ts} {level} worker={i%32:02d} tenant=T{i%380:04d} route={route} latency_ms={8+(i*13)%820} request={hashlib.md5(f"{day}-{i}".encode()).hexdigest()[:12]} {sentence(5,12)}\n')
    # Real rotated log encodings, including formats CMPCT should often leave alone.
    for q in rawfiles[:2]:
        b=q.read_bytes(); (p/(q.name+'.gz')).write_bytes(gzip.compress(b,compresslevel=6)); (p/(q.name+'.xz')).write_bytes(lzma.compress(b,preset=6))
    subprocess.run(['zstd','-q','-f','-8',str(rawfiles[2]),'-o',str(p/(rawfiles[2].name+'.zst'))],check=True)


def corpus_backups(root: Path) -> None:
    reseed('corpus_backups')
    p=root/'06_incremental_backups'; reset_dir(p)
    # Four snapshots share most content but include edits, deletes, and binary churn.
    seedfiles={}
    for i in range(180): seedfiles[f'doc_{i:03d}.txt']='\n'.join(sentence() for _ in range(45))
    for snap in range(4):
        d=p/f'snapshot_{snap}'; d.mkdir()
        for i,(name,body) in enumerate(seedfiles.items()):
            txt=body
            if i%17==snap%17: txt += '\nPATCH '+sentence(20,35)
            (d/name).write_text(txt,encoding='utf-8')
        for j in range(12):
            b=bytearray(randbytes(160*1024));
            if snap>0: b[:4096]=bytes([j+snap])*4096
            (d/f'blob_{j:02d}.bin').write_bytes(b)
    # Archive one snapshot as users often do, creating an already-compressed nested backup alongside raw snapshots.
    with zipfile.ZipFile(p/'snapshot_2.zip','w',compression=zipfile.ZIP_DEFLATED,compresslevel=6) as zf:
        for q in sorted((p/'snapshot_2').rglob('*')):
            if q.is_file(): zf.write(q,q.relative_to(p/'snapshot_2').as_posix())


def corpus_incompressible(root: Path) -> None:
    reseed('corpus_incompressible')
    p=root/'07_incompressible_and_encrypted_like'; reset_dir(p)
    (p/'encrypted_backup.bin').write_bytes(randbytes(8*1024*1024))
    for i in range(24): (p/f'chunk_{i:03d}.enc').write_bytes(randbytes(64*1024 + (i%7)*113))
    # Tiny random files are a metadata/expansion trap.
    tiny=p/'tokens'; tiny.mkdir()
    for i in range(1400): (tiny/f'{i:04d}.bin').write_bytes(randbytes(37 + i%233))


def corpus_tinyfiles(root: Path) -> None:
    reseed('corpus_tinyfiles')
    p=root/'08_many_tiny_files'; reset_dir(p)
    for i in range(5000):
        d=p/f'{i%90:02d}'/f'{(i//90)%20:02d}'; d.mkdir(parents=True,exist_ok=True)
        if i%7:
            data=(f'id={i}\nkind={i%19}\nstatus={"active" if i%3 else "idle"}\n'+sentence(5,14)+'\n').encode()
        else:
            data=randbytes(60+i%700)
        (d/f'item_{i:05d}.{ "json" if i%5 else "dat"}').write_bytes(data)


def corpus_ml(root: Path) -> None:
    reseed('corpus_ml')
    p=root/'09_ml_artifacts'; reset_dir(p)
    # Quantized-like weights: high entropy bytes plus repeated scale tables and structured metadata.
    weights=RNG.integers(0,256,size=12*1024*1024,dtype=np.uint8)
    # Introduce a realistic modest amount of zero/repeated blocks, not an absurdly easy tensor.
    for off in range(0,len(weights),512*1024):
        if (off//(512*1024))%7==0: weights[off:off+16384]=0
    (p/'model.q4.bin').write_bytes(weights.tobytes())
    scales=RNG.normal(0.03,0.008,size=600000).astype('<f2'); np.save(p/'scales.npy',scales)
    tokenizer={'model':'synthetic-bpe','vocab':{f'token_{i}':i for i in range(52000)},'merges':[f'token_{i} token_{(i*7)%52000}' for i in range(70000)]}
    (p/'tokenizer.json').write_text(json.dumps(tokenizer,separators=(',',':')),encoding='utf-8')
    (p/'training.log').write_text('\n'.join(f'step={i} loss={5/(1+i/4000):.6f} lr={2e-4*(1-i/30000):.9f} grad={0.4+0.1*math.sin(i/311):.5f}' for i in range(30000)),encoding='utf-8')


def corpus_disk(root: Path) -> None:
    reseed('corpus_disk')
    p=root/'10_large_mixed_binary'; reset_dir(p)
    q=p/'vm_disk_like.img'
    with q.open('wb') as f:
        for i in range(64):
            # 512 KiB blocks alternate among zeroed free space, structured pages, repeated OS-like pages,
            # and random encrypted/compressed-looking extents.  The changing entropy stresses chunk policy.
            mode=i%8
            if mode in (0,1): b=bytes(512*1024)
            elif mode in (2,3):
                page=(b'SQLITE_PAGE_SIM\0'+bytes([mode])*15+randbytes(160)+b'ROW|tenant=0042|status=active|')
                b=(page*((512*1024)//len(page)+1))[:512*1024]
            elif mode in (4,):
                page=(b'ELF_PAGE_SIM\0'+struct.pack('<I',i)+bytes(range(128)))
                b=(page*((512*1024)//len(page)+1))[:512*1024]
            else: b=randbytes(512*1024)
            f.write(b)


def build(root: Path) -> dict:
    reset_dir(root)
    builders=[corpus_source_repo,corpus_office,corpus_media,corpus_analytics,corpus_logs,corpus_backups,corpus_incompressible,corpus_tinyfiles,corpus_ml,corpus_disk]
    for fn in builders:
        print('building',fn.__name__,flush=True); fn(root)
    corpora=[]
    for d in sorted(x for x in root.iterdir() if x.is_dir()):
        files=[p for p in d.rglob('*') if p.is_file()]
        corpora.append({'name':d.name,'files':len(files),'logical_bytes':sum(p.stat().st_size for p in files),'tree_sha256':tree_hash(d)})
    manifest={'schema':'cmpct-neutral-hostile-corpus-v1','seed':SEED,'generated_utc':datetime.now(timezone.utc).isoformat(),
        'reproducibility_note':'Payload generation is deterministic per workload for a fixed dependency/toolchain set; container/media bytes can vary across external encoder versions, so benchmark records must pin tool versions.','contract':{
        'purpose':'Broad synthetic-but-valid regression corpus; not tuned to CMPCT.',
        'promotion_rule':'Report every workload, macro average, byte-weighted aggregate, worst-case expansion, creation/extraction, and selective-read cost. No corpus may be silently dropped.',
        'integrity':'Every workload has deterministic tree SHA-256.'},'corpora':corpora}
    (root/'MANIFEST.json').write_text(json.dumps(manifest,indent=2),encoding='utf-8')
    return manifest


if __name__=='__main__':
    ap=argparse.ArgumentParser(); ap.add_argument('--root',type=Path,default=ROOT); args=ap.parse_args()
    m=build(args.root); print(json.dumps(m,indent=2))
